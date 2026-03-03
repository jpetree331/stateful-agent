"""
Document reading tool: PDF and plain-text files.

Handles PDF extraction via pypdf and falls back to plain text reading for
everything else (.txt, .md, .csv, .log, .json, etc.).

Also provides extract_text_from_document_bytes() for Discord/Telegram attachments
(PDF, DOCX, PPTX, TXT, MD).

Dependencies:
  pypdf>=4.0, python-docx>=1.0, python-pptx>=0.6
"""
from __future__ import annotations

import io
from pathlib import Path

from langchain_core.tools import tool

_MAX_CHARS = 80_000  # generous limit; PDFs can be large

# File extensions supported for Discord/Telegram attachment extraction
DOCUMENT_EXTENSIONS = {".pdf", ".txt", ".pptx", ".docx", ".md"}


def extract_text_from_document_bytes(data: bytes, filename: str) -> str:
    """
    Extract text from document bytes (e.g. Discord attachment). Returns extracted
    text or error message. Handles PDF, DOCX, PPTX, TXT, MD.
    """
    fn = (filename or "").lower()
    ext = "." + fn.split(".")[-1] if "." in fn else ""
    if ext not in DOCUMENT_EXTENSIONS:
        return f"Unsupported format: {filename}"

    try:
        if ext == ".pdf":
            return _read_pdf_bytes(data)
        if ext == ".docx":
            return _read_docx_bytes(data)
        if ext == ".pptx":
            return _read_pptx_bytes(data)
        # .txt, .md and any other text
        return data.decode("utf-8", errors="replace")
    except Exception as e:
        return f"Error extracting text from {filename}: {e}"


def _normalize_pdf_text(text: str) -> str:
    """
    Fix fragmented PDF text where pypdf puts each word on its own line,
    separated by single or double newlines (word\\nword or word\\n\\nword).

    Approach: scan token by token. A "token" is any non-empty run of text
    between newlines. Short tokens (< 60 chars) are joined with spaces;
    long tokens (>= 60 chars) are treated as real paragraph breaks before them.
    """
    import re

    # Normalize line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    # Split on any sequence of newlines to get all tokens
    tokens = [t.strip() for t in re.split(r"\n+", text) if t.strip()]

    result = []
    buf = []

    for token in tokens:
        # Page markers always flush and stand alone
        if token.startswith("--- ") and token.endswith(" ---"):
            if buf:
                result.append(" ".join(buf))
                buf = []
            result.append(token)
        elif len(token) < 60:
            # Short token — likely a word fragment, accumulate
            buf.append(token)
        else:
            # Long token — real sentence/paragraph content
            if buf:
                # Check if this long token continues the buffer or starts fresh.
                # If buffer has only short fragments, join them with this token.
                result.append(" ".join(buf) + " " + token)
                buf = []
            else:
                result.append(token)

    if buf:
        result.append(" ".join(buf))

    return "\n\n".join(result)


def _read_pdf_bytes(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "Error: pypdf is not installed. Run: pip install pypdf"
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            normalized = _normalize_pdf_text(text.strip())
            pages.append(f"--- Page {i} ---\n{normalized}")
    return "\n\n".join(pages) if pages else "(No extractable text found in PDF)"


def _read_docx_bytes(data: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        return "Error: python-docx is not installed. Run: pip install python-docx"
    doc = Document(io.BytesIO(data))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_pptx_bytes(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx is not installed. Run: pip install python-pptx"
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        if slide_text:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
    return "\n\n".join(parts) if parts else "(No extractable text found in PPTX)"


def _read_pdf(path: Path) -> str:
    """Extract text from a PDF using pypdf."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return "Error: pypdf is not installed. Run: pip install pypdf"

    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            normalized = _normalize_pdf_text(text.strip())
            pages.append(f"--- Page {i} ---\n{normalized}")
    return "\n\n".join(pages) if pages else "(No extractable text found in PDF)"


@tool
def read_document(path: str) -> str:
    """
    Read and return the text content of a document — PDF or any plain-text file.

    For PDFs, extracts text page by page. For all other files (.txt, .md, .csv,
    .json, .log, code files, etc.) reads as plain text with UTF-8 encoding.
    Large documents are truncated at 80,000 characters with a notice.

    Use this instead of read_file when working with PDFs. For regular text files
    either tool works, but read_document is preferred for documents you intend to
    read and summarise.

    Args:
        path: Absolute or relative path to the document (~ is expanded).
    """
    try:
        p = Path(path).expanduser().resolve()
        if not p.exists():
            return f"Error: file not found — {p}"
        if not p.is_file():
            return f"Error: not a file — {p}"

        size_kb = p.stat().st_size / 1024
        header = f"[{p.name} | {size_kb:.1f} KB]\n\n"

        if p.suffix.lower() == ".pdf":
            content = _read_pdf(p)
        else:
            content = p.read_text(encoding="utf-8", errors="replace")

        truncated = ""
        if len(content) > _MAX_CHARS:
            content = content[:_MAX_CHARS]
            truncated = f"\n\n[... truncated at {_MAX_CHARS:,} chars — {size_kb:.0f} KB total]"

        return header + content + truncated

    except Exception as e:
        return f"Error reading {path}: {e}"
